import datetime
import itertools

import pymysql.cursors
import pymysql.cursors
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt

from data.model.Calendar_deal import Calendar_deal

file_date = datetime.datetime.today().strftime('%#m.%d.%Y')
today = datetime.date.today()
today_header = datetime.datetime.today().strftime('%m.%d.%Y')
providers = 'verizon', 'att', 'tmobile', 'sprint', 'xfinity'
Categories_title = 'BOGOF', 'Smartphone Other', 'Tablet', 'data Plan/Network', 'Trade-in', 'Switcher'
Categories_ref = 'bogo', 'smartphone other', 'tablet', 'data plan/network', 'trade-in'
day_of_week = datetime.timedelta(today.weekday())
calendar_start = today - (day_of_week - datetime.timedelta(days=4)) - datetime.timedelta(days=7*9)


def gen_col_dates(today):
    day_of_week = datetime.timedelta(today.weekday())
    factor = day_of_week - datetime.timedelta(days=4)
    col11 = (today + factor + datetime.timedelta(days=7)).strftime('%#m/%d')
    col10 = (today + factor).strftime('%#m/%d')
    col9 = (today + factor - datetime.timedelta(days=7)).strftime('%#m/%d')
    col8 = (today + factor - datetime.timedelta(days=7*2)).strftime('%#m/%d')
    col7 = (today + factor - datetime.timedelta(days=7*3)).strftime('%#m/%d')
    col6 = (today + factor - datetime.timedelta(days=7*4)).strftime('%#m/%d')
    col5 = (today + factor - datetime.timedelta(days=7*5)).strftime('%#m/%d')
    col4 = (today + factor - datetime.timedelta(days=7*6)).strftime('%#m/%d')
    col3 = (today + factor - datetime.timedelta(days=7*7)).strftime('%#m/%d')
    col2 = (today + factor - datetime.timedelta(days=7*8)).strftime('%#m/%d')
    col1 = (today + factor - datetime.timedelta(days=7*9)).strftime('%#m/%d')
    return [col1, col2, col3, col4, col5, col6, col7, col8, col9, col10, col11]


def gen_col_month(today):
    day_of_week = datetime.timedelta(today.weekday())
    factor = datetime.timedelta(days=4) - day_of_week
    month11 = (today + factor + datetime.timedelta(days=7)).strftime('%B')
    month10 = (today + factor).strftime('%B')
    month9 = (today + factor - datetime.timedelta(days=7)).strftime('%B')
    month8 = (today + factor - datetime.timedelta(days=7*2)).strftime('%B')
    month7 = (today + factor - datetime.timedelta(days=7*3)).strftime('%B')
    month6 = (today + factor - datetime.timedelta(days=7*4)).strftime('%B')
    month5 = (today + factor - datetime.timedelta(days=7*5)).strftime('%B')
    month4 = (today + factor - datetime.timedelta(days=7*6)).strftime('%B')
    month3 = (today + factor - datetime.timedelta(days=7*7)).strftime('%B')
    month2 = (today + factor - datetime.timedelta(days=7*8)).strftime('%B')
    month1 = (today + factor - datetime.timedelta(days=7*9)).strftime('%B')
    return month1[:3], month2[:3], month3[:3], month4[:3], month5[:3], month6[:3], month7[:3], month8[:3], month9[:3], month10[:3], month11[:3]


# determine how many columns have the same month
def month_cols_and_labels(months):
    row_count = 1
    month_row_amt = []
    month_row_labels = []
    for i in range(10):
        if months[i] == months[i + 1]:
            row_count += 1
            if i == 9:
                month_row_amt.append(row_count)
                month_row_labels.append(months[i])
        else:
            month_row_amt.append(row_count)
            month_row_labels.append(months[i])
            row_count = 1
            if i == 9:
                month_row_amt.append(row_count)
                month_row_labels.append(months[i + 1])
    return month_row_amt, month_row_labels


def get_deals(provider, category):
    connection = pymysql.connect(host='localhost',
                                 user='root',
                                 port=3306,
                                 password='123456',
                                 charset='utf8')

    query = "SELECT provider, category, deal_id, devices, promotion_details, promotion_summary, url, status, " \
            "modified_summary, date, homepage, start_date FROM historical_promotions WHERE provider = %s AND category = %s;"
    args = provider, category
    try:
        cursor = connection.cursor()
        cursor.execute('USE promotions')
        cursor.execute(query, args)
        deal_objs = []
        for deal in cursor.fetchall():
            deal_obj = Calendar_deal(deal[0], deal[1], deal[2], deal[3], deal[4], deal[5], deal[6], deal[7], deal[8],
                                     deal[9], deal[10], deal[11])
            deal_objs.append(deal_obj)
        cursor.close()
        return deal_objs
    finally:
        connection.commit()
        connection.close()


def get_flagship_deals(provider):
    connection = pymysql.connect(host='localhost',
                                 user='root',
                                 port=3306,
                                 password='123456',
                                 charset='utf8')

    query = "SELECT provider, category, deal_id, devices, promotion_details, promotion_summary, url, status, " \
            "modified_summary, date, homepage, start_date FROM historical_promotions WHERE provider = %s;"
    args = provider
    try:
        cursor = connection.cursor()
        cursor.execute('USE pricing')
        cursor.execute(query, args)
        deal_objs = []
        for deal in cursor.fetchall():
            deal_obj = Calendar_deal(deal[0], deal[1], deal[2], deal[3], deal[4], deal[5], deal[6], deal[7], deal[8],
                                     deal[9], deal[10], deal[11])
            if (deal_obj.end_date_ref_flagship - calendar_start).days > 5:
                deal_objs.append(deal_obj)
        cursor.close()
        return deal_objs
    finally:
        connection.commit()
        connection.close()


def generate_heading(col_dates, col_months, slide):
    # generate table for week date labels
    shapes = slide.shapes
    rows = 1
    cols = 11
    top = Inches(.798)
    left = Inches(1.25)
    width = Inches(11.575)
    height = Inches(0.212)
    table2 = shapes.add_table(rows, cols, left, top, width, height).table

    # set table color by column
    for i in range(11):
        fill2 = table2.cell(0, i).fill
        fill2.solid()
        fill2.fore_color.rgb = RGBColor(228, 229, 227)
        fill2.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_2

    # set font color and add text runs by column
    for i in range(11):
        text_frame = table2.cell(0, i).text_frame
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = col_dates[i]
        run.font.name = 'NeueHaasGroteskText Std (Body)'
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0, 0, 0)

    month_row_amt, month_row_labels = month_cols_and_labels(col_months)

    # generate table for month labels
    rows = 1
    cols = len(month_row_amt)
    top = Inches(0.585)
    left = Inches(1.25)
    width = Inches(11.575)
    height = Inches(0.212)
    table1 = shapes.add_table(rows, cols, left, top, width, height).table

    for i in range(len(month_row_amt)):
        # by row
        table1.columns[i].width = Inches(1.05 * month_row_amt[i])
        fill1 = table1.cell(0, i).fill
        fill1.solid()
        fill1.fore_color.rgb = RGBColor(173, 175, 175)
        text_frame = table1.cell(0, i).text_frame
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = month_row_labels[i]
        run.font.italic = True
        run.font.bold = True
        run.font.name = 'NeueHaasGroteskText Std (Body)'
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0, 0, 0)


def add_boxes_to_calendar(slide, dictionary):
    # edit start & end dates of duplicate deal_ids
    for x in range(5):
        for a, b in itertools.combinations(dictionary[providers[x]], 2):
            if a.deal_id == b.deal_id:
                if a.end_date_cal >= b.end_date_cal:
                    a.start_date_cal = b.end_date_cal
                    a.start_date = b.end_date
                    a.start_date_ref = b.end_date_ref
                else:
                    b.start_date_cal = a.end_date_cal
                    b.start_date = a.end_date
                    b.start_date_ref = a.end_date_ref
        # set rows to be empty list
        rows = {}
        # maximum number of rows is the number of deals, so make a row for each deal
        for deal_count in range(len(dictionary[providers[x]])):
            rows[deal_count] = []
        # condense deals that do not overlap into the same row
        for deal in sorted(dictionary[providers[x]], key=lambda object: object.start_date_cal):
            for row_number in rows:
                if rows[row_number]:
                    if deal.start_date_cal >= sorted(rows[row_number], key=lambda object: object.end_date_cal)[-1].end_date_cal:
                        rows[row_number].append(deal)
                        break
                    else:
                        continue
                else:
                    rows[row_number].append(deal)
                    break
        # remove empty rows
        for row_number in range(len(rows)):
            if not rows[row_number]:
                del rows[row_number]
        # add text boxes for each provider
        levels = len(rows)
        for row_number in range(levels):
            for deal in rows[row_number]:
                top = Inches(1.06 + (1.178 * x) + (row_number * 1.178 / levels))
                if (deal.start_date_ref - calendar_start).days <= 0:
                    left_inches = 1.25
                else:
                    left_inches = 1.25 + (deal.start_date_ref - calendar_start).days * (12.2 / 76)
                left = Inches(left_inches)
                if left_inches == 1.25:
                    width_inches = (deal.end_date_ref - calendar_start).days * (12.2 / 76)
                else:
                    width_inches = (deal.end_date_ref - deal.start_date_ref).days * (12.2 / 76)
                height = Inches(1.06 / levels)
                width = Inches(width_inches)
                text_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
                fill = text_box.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(242, 242, 242)
                line = text_box.line
                line.color.rgb = RGBColor(0, 0, 0)
                tf = text_box.text_frame
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = deal.promotion_summary + " (" + str(deal.start_date) + "-" + deal.end_date + ")"
                run.font.name = 'NeueHaasGroteskText Std (Body)'
                if deal.iconic == 'yes':
                    run.font.color.rgb = RGBColor(0, 112, 192)
                elif deal.iconic == 'no':
                    run.font.color.rgb = RGBColor(0, 0, 0)
                run.font.bold = False
                if deal.status != 'discontinued':
                    run.font.bold = True
                if width_inches < 3.5:
                    run.font.size = Pt(7)
                else:
                    run.font.size = Pt(8)


def generate_PowerPoint():
    prs = Presentation(
        r"C:\Users\Amanda Friedman\PycharmProjects\DailyPromotionsAndPricing\Templates\Daily-PowerPoint-Template-Calendar-Xfinity.pptx")

    # title text
    slide = prs.slides[0]
    left = Inches(0.42)
    top = Inches(0.1)
    width = Inches(6.0)
    height = Inches(1.2)
    title_text = slide.shapes.add_textbox(left, top, width, height)
    tf = title_text.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Competitive Overview: Smartphone Promotions (Flagship Handsets)"
    run.font.bold = True
    run.font.name = 'NeueHaasGroteskDisp Std (Headings)'
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(0, 0, 0)

    # sub title text with date
    left = Inches(0.42)
    top = Inches(0.32)
    width = Inches(6.0)
    height = Inches(.4)
    sub_title_text = slide.shapes.add_textbox(left, top, width, height)
    tf = sub_title_text.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "10 Week Trailing Calendar (updated through (" + today_header + ")"
    run.font.bold = True
    run.font.name = 'NeueHaasGroteskDisp Std (Headings)'
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(255, 0, 0)

    # add text to top corner
    left = Inches(11.6)
    top = Inches(0.1)
    width = Inches(1)
    height = Inches(0.5)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = 'Bold = currently in market'
    run.font.bold = True
    run.font.name = 'NeueHaasGroteskText Std (Body)'
    run.font.size = Pt(7.5)
    run.font.color.rgb = RGBColor(0, 0, 0)
    run = p.add_run()
    run.text = '\nBlue = iconic offers'
    run.font.bold = True
    run.font.name = 'NeueHaasGroteskText Std (Body)'
    run.font.size = Pt(7.5)
    run.font.color.rgb = RGBColor(0, 112, 192)

    col_dates = gen_col_dates(today)
    col_months = gen_col_month(today)

    generate_heading(col_dates, col_months, slide)

    add_boxes_to_calendar(slide, flagship_deal_dict)

    slide = prs.slides[1]
    left = Inches(0.42)
    top = Inches(0.1)
    width = Inches(6.0)
    height = Inches(1.2)
    title_text = slide.shapes.add_textbox(left, top, width, height)
    tf = title_text.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Competitive Overview: Smartphone Promotions (Entry Handsets)"
    run.font.bold = True
    run.font.name = 'NeueHaasGroteskDisp Std (Headings)'
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(0, 0, 0)

    # sub title text with date
    left = Inches(0.42)
    top = Inches(0.32)
    width = Inches(6.0)
    height = Inches(.4)
    sub_title_text = slide.shapes.add_textbox(left, top, width, height)
    tf = sub_title_text.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "10 Week Trailing Calendar (updated through (" + today_header + ")"
    run.font.bold = True
    run.font.name = 'NeueHaasGroteskDisp Std (Headings)'
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(255, 0, 0)

    # add text to top corner
    left = Inches(11.6)
    top = Inches(0.1)
    width = Inches(1)
    height = Inches(0.5)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = 'Bold = currently in market'
    run.font.bold = True
    run.font.name = 'NeueHaasGroteskText Std (Body)'
    run.font.size = Pt(7.5)
    run.font.color.rgb = RGBColor(0, 0, 0)

    add_boxes_to_calendar(slide, budget_deal_dict)

    generate_heading(col_dates, col_months, slide)

    prs.save(r"C:\Users\Amanda Friedman\PycharmProjects\DailyPromotionsAndPricing\Calendar Slides\Device Promotions Competitive Calendar (" + file_date + ") Tarifica Version- draft.pptx")
    print("PowerPoint generated.")


deal_dict = {}
flagship_deal_dict = {}
budget_deal_dict = {}
for x in range(5):
    deal_dict[providers[x]] = get_flagship_deals(providers[x])
    flagship_deal_dict[providers[x]] = []
    budget_deal_dict[providers[x]] = []
    for deal in deal_dict[providers[x]]:
        if deal.flagship == 'yes':
            flagship_deal_dict[providers[x]].append(deal)
        elif deal.budget == 'yes':
            budget_deal_dict[providers[x]].append(deal)

generate_PowerPoint()