import datetime
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR

def col_dates(today):
    day_of_week = datetime.timedelta(today.weekday())
    col11 = (today - day_of_week + datetime.timedelta(days=7)).strftime('%#m/%d')
    col10 = (today - day_of_week).strftime('%#m/%d')
    col9 = (today - day_of_week - datetime.timedelta(days=7)).strftime('%#m/%d')
    col8 = (today - day_of_week - datetime.timedelta(days=7*2)).strftime('%#m/%d')
    col7 = (today - day_of_week - datetime.timedelta(days=7*3)).strftime('%#m/%d')
    col6 = (today - day_of_week - datetime.timedelta(days=7*4)).strftime('%#m/%d')
    col5 = (today - day_of_week - datetime.timedelta(days=7*5)).strftime('%#m/%d')
    col4 = (today - day_of_week - datetime.timedelta(days=7*6)).strftime('%#m/%d')
    col3 = (today - day_of_week - datetime.timedelta(days=7*7)).strftime('%#m/%d')
    col2 = (today - day_of_week - datetime.timedelta(days=7*8)).strftime('%#m/%d')
    col1 = (today - day_of_week - datetime.timedelta(days=7*9)).strftime('%#m/%d')
    return [col1, col2, col3, col4, col5, col6, col7, col8, col9, col10, col11]

def col_month(today):
    day_of_week = datetime.timedelta(today.weekday())
    month11 = (today - day_of_week + datetime.timedelta(days=7)).strftime('%B')
    month10 = (today - day_of_week).strftime('%B')
    month9 = (today - day_of_week - datetime.timedelta(days=7)).strftime('%B')
    month8 = (today - day_of_week - datetime.timedelta(days=7*2)).strftime('%B')
    month7 = (today - day_of_week - datetime.timedelta(days=7*3)).strftime('%B')
    month6 = (today - day_of_week - datetime.timedelta(days=7*4)).strftime('%B')
    month5 = (today - day_of_week - datetime.timedelta(days=7*5)).strftime('%B')
    month4 = (today - day_of_week - datetime.timedelta(days=7*6)).strftime('%B')
    month3 = (today - day_of_week - datetime.timedelta(days=7*7)).strftime('%B')
    month2 = (today - day_of_week - datetime.timedelta(days=7*8)).strftime('%B')
    month1 = (today - day_of_week - datetime.timedelta(days=7*9)).strftime('%B')
    return month1[:3], month2[:3], month3[:3], month4[:3], month5[:3], month6[:3], month7[:3], month8[:3], month9[:3], month10[:3], month11[:3]

# determine how many columns have the same month
def month_cols_and_labels(months):
    row_count = 1
    month_row_amt = []
    month_row_labels = []
    for i in range(10):
        if months[i] == months[i + 1]:
            row_count += 1
            if i == 9:
                month_row_amt.append(row_count)
                month_row_labels.append(months[i])
        else:
            month_row_amt.append(row_count)
            month_row_labels.append(months[i])
            row_count = 1
            if i == 9:
                month_row_amt.append(row_count)
                month_row_labels.append(months[i + 1])
    return month_row_amt, month_row_labels

def generate_calendar_heading(calendar_col_dates, calendar_col_months, slide):
    # generate table for week date labels
    shapes = slide.shapes
    rows = 1
    cols = 11
    top = Inches(1.01)
    left = Inches(1.25)
    width = Inches(11.575)
    height = Inches(0.425)
    table2 = shapes.add_table(rows, cols, left, top, width, height).table

    # set table color by column
    for i in range(11):
        fill2 = table2.cell(0, i).fill
        fill2.solid()
        fill2.fore_color.rgb = RGBColor(228, 229, 227)
        fill2.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_2

    # set font color and add text runs by column
    for i in range(11):
        text_frame = table2.cell(0, i).text_frame
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = calendar_col_dates[i]
        run.font.name = 'NeueHaasGroteskText Std (Body)'
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0, 0, 0)

    month_row_amt, month_row_labels = month_cols_and_labels(calendar_col_months)

    # generate table for month labels
    rows = 1
    cols = len(month_row_amt)
    top = Inches(0.585)
    left = Inches(1.25)
    width = Inches(11.575)
    height = Inches(0.425)
    table1 = shapes.add_table(rows, cols, left, top, width, height).table

    for i in range(len(month_row_amt)):
        # by row
        table1.columns[i].width = Inches(1.05 * month_row_amt[i])
        fill1 = table1.cell(0, i).fill
        fill1.solid()
        fill1.fore_color.rgb = RGBColor(173, 175, 175)
        text_frame = table1.cell(0, i).text_frame
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = month_row_labels[i]
        run.font.italic = True
        run.font.bold = True
        run.font.name = 'NeueHaasGroteskText Std (Body)'
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0, 0, 0)