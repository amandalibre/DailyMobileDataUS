import csv
import datetime
import itertools
import re

from daily_promotions_and_pricing.Deal_Errors import deals_errors
from docx import Document
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt

from data.database.Database_Methods import check_duplicates, edit_yesterday, get_calendar_deals, remove_yesterday, \
    add_to_database
from data.model.Deal import Deal
from daily_promotions_and_pricing.Calendar_Methods import col_month, col_dates, month_cols_and_labels
from daily_promotions_and_pricing.Device_Slides import generate_Device_Slides

deals_by_provider = {}
NUM_PROVIDERS = 6
provider_names = ['verizon', 'att', 'tmobile', 'sprint', 'metropcs', 'cricket']
provider_names_email = ['Verizon', 'AT&T', 'T-Mobile', 'Sprint', 'MetroPCS', 'Cricket']
today = datetime.date.today()
day_of_week = datetime.datetime.today().weekday()
today_header = datetime.datetime.today().strftime('%m/%d/%Y')
today_cover = datetime.datetime.today().strftime('%B' + ' ' + '%d' + ', ' + '%Y')
today_filename = datetime.datetime.today().strftime('%m.%d.%Y')
today_deal_id = datetime.datetime.today().strftime('%Y%m%d')
calendar_col_dates = []
calendar_providers = 'verizon', 'att', 'tmobile', 'sprint'
Categories_title = 'BOGOF', 'Smartphone Other', 'Tablet', 'data Plan/Network', 'Trade-in', 'Switcher'
Categories_ref = 'bogo', 'smartphone other', 'tablet', 'data plan/network', 'trade-in'
font_color_provider = RGBColor(192, 0, 0), RGBColor(0, 112, 192), RGBColor(112, 48, 160), RGBColor(0, 102, 0)
calendar_start = today - datetime.timedelta(today.weekday()) - datetime.timedelta(days=7*9)

approved_device_list = []
def get_approved_device_list():
    provider_file = r'C:\Users\Amanda Friedman\PycharmProjects\DailyPromotionsAndPricing\Master Device List\Device List.csv'
    with open(provider_file, 'r') as f:
        reader = csv.reader(f)
        next(reader)
        for row in reader:
            device = row[0]
            device = device.strip()
            approved_device_list.append(device)
    return approved_device_list


def get_day_before(today):
    if day_of_week == 0:
        day_before = today - datetime.timedelta(days=2)
    else:
        day_before = today - datetime.timedelta(days=1)
    return day_before

def get_deals(provider_name):
    deals = []
    provider_file = r'C:/Users/Amanda Friedman/Documents/Verizon/Daily Tracking Files/VerizonPromotionsData/' + provider_name + '-' + str(today) + '.csv'
    with open(provider_file, 'r') as f:
        reader = csv.reader(f)
        next(reader)
        for row in reader:
            provider, category, deal_id, devices, promotion_details, promotion_summary, url, status, modified_summary, date, homepage = row
            filename_provider, filename_date = provider_name, str(today)
            deal = Deal(provider, category, deal_id, devices, promotion_details, promotion_summary, url, status, modified_summary, date, homepage, filename_provider, filename_date)
            deals.append(deal)
    print("-> " + provider_name + " file read")
    return deals

def generate_cover_email(deals_by_provider):
    document = Document()
    p = document.add_paragraph()
    p_format = p.paragraph_format
    p_format.space_before = Pt(0)
    p_format.space_after = Pt(0)
    # p_format.line_spacing = Pt(0)
    run = p.add_run()
    font = run.font
    font.name = "Calibri (Body)"
    font.size = Pt(11)
    run.text = " \n"
    run = p.add_run()
    run.text = "Dear Verizon Team," + "\n" + "\n"
    run = p.add_run()
    run.text = "Attached is today’s pricing and promotions PowerPoint."
    for i in range(6):
        p = document.add_paragraph()
        p_format = p.paragraph_format
        p_format.space_before = Pt(0)
        p_format.space_after = Pt(0)
        # p_format.line_spacing = Pt(0)
        run = p.add_run()
        run.text = "\n" + provider_names_email[i]
        font = run.font
        font.bold = True
        deal_count = len(deals_by_provider[provider_names[i]])
        loop_count = 0
        changes = 0
        for promo in deals_by_provider[provider_names[i]]:
            loop_count = loop_count + 1
            if promo.status == "new":
                changes = changes + 1
                p = document.add_paragraph()
                p_format = p.paragraph_format
                p_format.space_before = Pt(0)
                p_format.space_after = Pt(0)
                # p_format.line_spacing = Pt(0)
                p.style = "List Bullet"
                p_format = p.paragraph_format
                p_format.left_indent = Inches(0.5)
                run = p.add_run()
                run.text = ''
                run.text = "New Offer: " + promo.promotion_details
                font = run.font
                font.bold = False
        for promo in deals_by_provider[provider_names[i]]:
            loop_count = loop_count + 1
            if promo.status == "modified":
                changes = changes + 1
                p = document.add_paragraph()
                p_format = p.paragraph_format
                p_format.space_before = Pt(0)
                p_format.space_after = Pt(0)
                # p_format.line_spacing = Pt(0)
                p.style = "List Bullet"
                p_format = p.paragraph_format
                p_format.left_indent = Inches(0.5)
                run = p.add_run()
                run.text = ''
                run.text = "Modified Offer: " + promo.promotion_details + " (" + promo.modified_summary + ")"
                font = run.font
                font.bold = False
        for promo in deals_by_provider[provider_names[i]]:
            loop_count = loop_count + 1
            if promo.status == "discontinued":
                changes = changes + 1
                p = document.add_paragraph()
                p_format = p.paragraph_format
                p_format.space_before = Pt(0)
                p_format.space_after = Pt(0)
                # p_format.line_spacing = Pt(0)
                p.style = "List Bullet"
                p_format = p.paragraph_format
                p_format.left_indent = Inches(0.5)
                run = p.add_run()
                run.text = ''
                run.text = "Discontinued Offer: " + promo.promotion_details
                font = run.font
                font.bold = False
        if loop_count == (deal_count*3) and changes == 0:
            run = p.add_run()
            run.text = ''
            run.text = " -- No changes"
            font = run.font
            font.bold = False
    p = document.add_paragraph()
    run = p.add_run()
    run.text = "\n" + "Best regards,"
    document.save(r"C:\Users\Amanda Friedman\PycharmProjects\DailyPromotionsAndPricing\Daily PowerPoint\Cover Email (" + str(today_filename) + ").docx")
    print("Cover email generated.")

def generate_promo_text(p, promo, color):
    promo_text = promo.promotion_details
    bold_words = re.findall(r'[$]\S*', promo_text)
    text_words = promo_text.split(" ")
    for word in text_words:
        if word in bold_words or word == "free" or word == "Free":
            run = p.add_run()
            run.text = word + " "
            font = run.font
            font.name = 'NeueHaasGroteskText Std (Body)'
            font.size = Pt(9)
            font.bold = True
            if color == "red":
                font.color.rgb = RGBColor(255, 0, 0)
            elif color == "blue":
                font.color.rgb = RGBColor(0, 176, 240)
            else:
                font.color.rgb = RGBColor(0, 0, 0)
        else:
            run = p.add_run()
            run.text = word + " "
            font = run.font
            font.name = 'NeueHaasGroteskText Std (Body)'
            font.size = Pt(9)
            font.bold = False
            if color == "red":
                font.color.rgb = RGBColor(255, 0, 0)
            elif color == "blue":
                font.color.rgb = RGBColor(0, 176, 240)
            else:
                font.color.rgb = RGBColor(0, 0, 0)
    run.text = text_words[-1] + " (" + promo.deal_id[6:8] + "/" + promo.deal_id[8:10] + "/" + promo.deal_id[4:6] + ")\n" + "\n"
    font = run.font
    font.bold = False

def generate_PowerPoint(deals_by_provider):
    prs = Presentation(
        r"C:\Users\Amanda Friedman\PycharmProjects\DailyPromotionsAndPricing\Templates\Daily-PowerPoint-Template-Full.pptx")

    #slides [0: cover, 1: device chart, 6: calendar, 7: promotions, 13: end]
    slides = [None] * 13
    slide_titles = ['', '', '', '', '', '', '', 'BOGOF', 'Smartphone Other', 'Tablet', 'data Plan/Network', 'Trade-in', 'Switcher', '']
    category_slide_refs = ['', '', '', '', '', '', '', 'bogo', 'smartphone other', 'tablet', 'data plan/network', 'trade-in',
                    'switcher', '']

    #cover slide
    #add date
    slide0 = prs.slides[0]
    left = Inches(0.6)
    top = Inches(3.1)
    width = Inches(3)
    height = Inches(0.4)
    txBox = slide0.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(today_cover)
    run.font.bold = True
    run.font.name = 'NeueHaasGroteskText Std (Body)'
    run.font.size = Pt(18)

    # device chart slides [1-6]
    generate_Device_Slides(prs)

    # calendar slide [6]
    # title text
    slide6 = prs.slides[6]
    left = Inches(0.45)
    top = Inches(0.1)
    width = Inches(6.0)
    height = Inches(1.2)
    title_text = slide6.shapes.add_textbox(left, top, width, height)
    tf = title_text.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = Categories_title[0] + " Promotions"
    run.font.bold = True
    run.font.name = 'NeueHaasGroteskDisp Std (Body)'
    run.font.size = Pt(26)
    run.font.color.rgb = RGBColor(205, 4, 11)

    #add date to top corner
    slide6 = prs.slides[6]
    left = Inches(11.8)
    top = Inches(0.1)
    width = Inches(1)
    height = Inches(0.3)
    txBox = slide6.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "as of " + str(today_header)
    run.font.italic = True
    run.font.name = 'NeueHaasGroteskText Std (Body)'
    run.font.size = Pt(10)

    calendar_col_dates = col_dates(today)
    calendar_col_months = col_month(today)

    # generate table for week date labels
    shapes = slide6.shapes
    rows = 1
    cols = 11
    top = Inches(1.01)
    left = Inches(1.25)
    width = Inches(11.575)
    height = Inches(0.425)
    table2 = shapes.add_table(rows, cols, left, top, width, height).table

    # set table color by column
    for i in range(11):
        fill2 = table2.cell(0, i).fill
        fill2.solid()
        fill2.fore_color.rgb = RGBColor(211, 211, 211)
        fill2.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_2

    # set font color and add text runs by column
    for i in range(11):
        text_frame = table2.cell(0, i).text_frame
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = calendar_col_dates[i]
        run.font.name = 'NeueHaasGroteskText Std (Body)'
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0, 0, 0)

    month_row_amt, month_row_labels = month_cols_and_labels(calendar_col_months)

    # generate table for month labels
    rows = 1
    cols = len(month_row_amt)
    top = Inches(0.585)
    left = Inches(1.25)
    width = Inches(11.575)
    height = Inches(0.425)
    table1 = shapes.add_table(rows, cols, left, top, width, height).table

    for i in range(len(month_row_amt)):
        # by row
        table1.columns[i].width = Inches(1.05 * month_row_amt[i])
        fill1 = table1.cell(0, i).fill
        fill1.solid()
        fill1.fore_color.rgb = RGBColor(249, 178, 149)
        text_frame = table1.cell(0, i).text_frame
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = month_row_labels[i]
        run.font.italic = True
        run.font.bold = True
        run.font.name = 'NeueHaasGroteskText Std (Body)'
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(0, 0, 0)

    # get deals from historical_promotions
    calendar_deal_dict = {}
    for x in range(4):
        calendar_deal_dict[calendar_providers[x]] = get_calendar_deals(calendar_providers[x], 'bogo')

    # remove duplicate deal_ids (temporary fix for modified deals with same start dates but earlier end dates)
    for x in range(4):
        for a, b in itertools.combinations(calendar_deal_dict[calendar_providers[x]], 2):
            if a.deal_id == b.deal_id:
                if a.end_date_cal >= b.end_date_cal:
                    a.start_date_cal = b.end_date_cal
                    a.start_date = b.end_date
                    a.start_date_ref = b.end_date_ref
                else:
                    b.start_date_cal = a.end_date_cal
                    b.start_date = a.end_date
                    b.start_date_ref = a.end_date_ref
    # add text boxes for promotions
    for x in range(4):
        boxes = len(calendar_deal_dict[calendar_providers[x]])
        y = 0
        for deal in calendar_deal_dict[calendar_providers[x]]:
            top = Inches(1.435 + (1.369 * x) + (y * 1.35 / boxes))
            if (deal.end_date_ref - calendar_start).days <= 5:
                continue
            if (deal.start_date_ref - calendar_start).days <= 0:
                left_inches = 1.25
            else:
                left_inches = 1.25 + (deal.start_date_ref - calendar_start).days * (11.575 / 76)
            left = Inches(left_inches)
            if left_inches == 1.25:
                width_inches = (deal.end_date_ref - calendar_start).days * (11.575 / 76)
            else:
                width_inches = (deal.end_date_ref - deal.start_date_ref).days * (11.575 / 76)
            width = Inches(width_inches)
            height = Inches(1.3426 / boxes)
            text_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
            fill = text_box.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(221, 221, 221)
            line = text_box.line
            line.color.rgb = RGBColor(0, 0, 0)
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = deal.promotion_summary + " (" + str(deal.start_date) + "-" + deal.end_date + ")"
            run.font.color.rgb = font_color_provider[x]
            run.font.name = 'NeueHaasGroteskText Std (Body)'
            run.font.size = Pt(10)
            if width_inches < 3:
                run.font.size = Pt(8)
            if width_inches < 2:
                run.font.size = Pt(7.5)
            run.font.bold = True
            y +=1

    # add line for current day
    left = Inches(10.84 + (11.575 / 76)*datetime.datetime.today().weekday())
    top = Inches(0.585)
    width = Inches(.015)
    height = Inches(6.28)
    slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

    # add box for current day
    left = Inches(10.47 + (11.575 / 76)*datetime.datetime.today().weekday())
    top = Inches(0.185)
    width = Inches(0.75)
    height = Inches(0.7)
    today_arrow = slide6.shapes.add_shape(MSO_SHAPE.DOWN_ARROW_CALLOUT, left, top, width, height)
    text_frame = today_arrow.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = 'TODAY' + '\n' + datetime.datetime.today().strftime('%m/%d')
    run.font.italic = False
    run.font.bold = False
    run.font.name = 'NeueHaasGroteskText Std (Body)'
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0, 0, 0)

    #promotions slides [7-13]
    #by slide
    for x in range(7, 13):
        slides[x] = (prs.slides[x])

        #changes title to category name
        title_placeholder = slides[x].shapes.title
        title_placeholder.text = "Promotions: " + slide_titles[x]

        #add date to top corner
        left = Inches(11.8)
        top = Inches(0.1)
        width = Inches(1)
        height = Inches(0.3)
        txBox = slides[x].shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        daterun = p.add_run()
        daterun.text = str(today_header)
        daterun.font.italic = True
        daterun.font.name = 'NeueHaasGroteskText Std (Body)'
        daterun.font.size = Pt(10)

        #generate table
        shapes = slides[x].shapes
        rows = 1
        cols = 6
        top = Inches(1.4)
        left = Inches(0.65)
        width = Inches(12.0)
        height = Inches(5.4)
        table = shapes.add_table(rows, cols, left, top, width, height).table

        # set column width
        table.columns.width = Inches(2.0)

        # set table color
        # by column
        for i in range(6):
            fill = table.cell(0, i).fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(211, 211, 211)
            fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_2

        # set font color and add text runs
        # by column
        for i in range(6):
            text_frame = table.cell(0, i).text_frame
            p = text_frame.paragraphs[0]
            # by promo
            for promo in deals_by_provider[provider_names[i]]:
                if promo.provider == provider_names[i] and promo.category == category_slide_refs[x] and promo.status != "discontinued":
                    if promo.status == "new" or promo.status == "modified":
                        generate_promo_text(p, promo, "red")
                    else:
                        if promo.homepage == "yes":
                            generate_promo_text(p, promo, "blue")
                        else:
                            generate_promo_text(p, promo, "black")
    prs.save(r"C:\Users\Amanda Friedman\PycharmProjects\DailyPromotionsAndPricing\Daily PowerPoint\Competitive Pricing Landscape (" + str(today_filename) + ") - draft.pptx")
    print("PowerPoint generated.")

def generate_PowerPoint_and_cover_email():

    for x in range(NUM_PROVIDERS):
        deals_by_provider[provider_names[x]] = get_deals(provider_names[x])

    get_approved_device_list()
    deals_errors(deals_by_provider, approved_device_list)

    daily_count = 0
    hist_count = 0
    num_deals = 0
    for x in range(NUM_PROVIDERS):
        num_deals += len(deals_by_provider[provider_names[x]])
        for deal in deals_by_provider[provider_names[x]]:
            if check_duplicates("daily_promotions", deal.deal_id, deal.date_mysql) == False:
                daily_count += 1
                add_to_database("daily_promotions", deal.provider, deal.category, deal.deal_id, deal.devices,
                                deal.promotion_details,
                                deal.promotion_summary, deal.url, deal.status, deal.modified_summary, deal.date_mysql,
                                deal.homepage, deal.start_date_mysql)
            if check_duplicates("historical_promotions", deal.deal_id, deal.date_mysql) == False:
                hist_count += 1
                add_to_database("historical_promotions", deal.provider, deal.category, deal.deal_id, deal.devices,
                                deal.promotion_details,
                                deal.promotion_summary, deal.url, deal.status, deal.modified_summary, deal.date_mysql,
                                deal.homepage, deal.start_date_mysql)
            day_before = get_day_before(today)
            if deal.status == 'modified':
                edit_yesterday(day_before, deal.deal_id, deal.date_mysql)
            remove_yesterday(day_before, deal.deal_id)

    print(daily_count, "deals added to daily database &", num_deals - daily_count, "duplicates ignored.")
    print(hist_count, "deals added &", num_deals - hist_count, "duplicates ignored.")

    generate_cover_email(deals_by_provider)

    generate_PowerPoint(deals_by_provider)


generate_PowerPoint_and_cover_email()