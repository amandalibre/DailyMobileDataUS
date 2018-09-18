import datetime

from docx.shared import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR

from data.database.Database_Methods import get_postpaid_device_prices, get_prepaid_device_prices
import itertools

today_filename = datetime.datetime.today().strftime('%m.%d.%Y')
today = datetime.date.today()
postpaid_providers = ['verizon', 'att', 'tmobile', 'sprint', 'xfinity']
prepaid_providers = ['verizon', 'att', 'metropcs', 'cricket']
time_now = datetime.datetime.now().time()

def fix_spacing(string):
    string = string.replace(u'\xa0', u' ')
    string = string.strip()
    return string

def format_money(string):
    string = '{:.2f}'.format(string)
    return string

def format_device_names(string):
    device = string[0]
    device = device.title()
    device = device.replace('Iphone', 'iPhone')
    device = device.replace('Ipad', 'iPad')
    device = device.replace('Lg', 'LG')
    device = device.replace('Zte', 'ZTE')
    device = device.replace('Htc', 'HTC')
    device = device.replace('Fierce', 'FIERCE')
    device = device.replace('At&t', 'AT&T')
    device = device.replace('Xl', 'XL')
    device = device.replace('Gizmotab', 'GizmoTab')
    device = device.replace('Se', 'SE')
    device = device.replace('6S', '6s')
    device = device.replace('Xp', 'XP')
    if device.find('Google') == -1 and device.find('Pixel') != -1:
        device = device.replace('Pixel', 'Google Pixel')
    storage = string[1]
    return device + ' (' + storage + ' GB)'

def populate_cell(key, provider, chart, cell, run, column):
    if key in chart[provider]:
        if column == 1 or column == 4 or column == 7 or column == 10 or column == 13:
            if format_money(chart[provider][key].monthly_price) == '0.00':
                if chart == tab_chart_2:
                    run.text = 'NA'
                else:
                    run.text = 'Free'
            else:
                run.text = '$' + format_money(chart[provider][key].monthly_price)
        elif column == 2 or column == 5 or column == 8 or column == 11 or column == 14:
            run.text = '$' + format_money(chart[provider][key].retail_price)
        elif column == 3 or column == 6 or column == 9 or column == 12:
            if format_money(chart[provider][key].onetime_price) == '0.00':
                run.text = 'NA'
            else:
                run.text = '$' + format_money(chart[provider][key].onetime_price)
        run.font.bold = True
        run.font.size = Pt(11)
        cell.fill.solid()
        if provider == 'verizon':
            cell.fill.fore_color.rgb = RGBColor(246, 231, 231)
        elif provider == 'att':
            cell.fill.fore_color.rgb = RGBColor(153, 204, 255)
        elif provider == 'tmobile':
            cell.fill.fore_color.rgb = RGBColor(237, 194, 217)
        elif provider == 'sprint':
            cell.fill.fore_color.rgb = RGBColor(179, 218, 180)
        elif provider == 'xfinity':
            cell.fill.fore_color.rgb = RGBColor(204, 204, 255)
        if column == 1 or column == 4 or column == 7 or column == 10 or column == 13:
            if chart[provider][key].monthly_price_change == 'yes':
                run.font.color.rgb = RGBColor(255, 0, 0)
            else:
                run.font.color.rgb = RGBColor(109, 110, 113)
        elif column == 2 or column == 5 or column == 8 or column == 11 or column == 14:
            if chart[provider][key].retail_price_change == 'yes':
                run.font.color.rgb = RGBColor(255, 0, 0)
            else:
                run.font.color.rgb = RGBColor(109, 110, 113)
        elif column == 3 or column == 6 or column == 9 or column == 12:
            if chart[provider][key].onetime_price_change == 'yes':
                run.font.color.rgb = RGBColor(255, 0, 0)
            else:
                run.font.color.rgb = RGBColor(109, 110, 113)
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(191, 191, 191)
        run.text = ' '

def populate_cell_contract_ufc(key, provider, chart, cell, run):
    if key in chart[provider]:
        if format_money(chart[provider][key].contract_ufc) == '0.00':
            run.text = 'NA'
        else:
            run.text = '$' + format_money(chart[provider][key].contract_ufc)
        run.font.bold = True
        run.font.size = Pt(11)
        cell.fill.solid()
        if provider == 'verizon':
            cell.fill.fore_color.rgb = RGBColor(246, 231, 231)
        elif provider == 'att':
            cell.fill.fore_color.rgb = RGBColor(153, 204, 255)
        elif provider == 'tmobile':
            cell.fill.fore_color.rgb = RGBColor(237, 194, 217)
        elif provider == 'sprint':
            cell.fill.fore_color.rgb = RGBColor(179, 218, 180)
        elif provider == 'xfinity':
            cell.fill.fore_color.rgb = RGBColor(204, 204, 255)
        run.font.color.rgb = RGBColor(109, 110, 113)
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(191, 191, 191)
        run.text = ' '

def prepaid_populate_cell(key, provider, chart, cell, run):
    if key in chart[provider]:
        run.text = '$' + str(chart[provider][key].retail_price)
        run.font.bold = True
        run.font.size = Pt(11)
        cell.fill.solid()
        if provider == 'verizon':
            cell.fill.fore_color.rgb = RGBColor(246, 231, 231)
        elif provider == 'att':
            cell.fill.fore_color.rgb = RGBColor(153, 204, 255)
        elif provider == 'metropcs':
            cell.fill.fore_color.rgb = RGBColor(253, 229, 161)
        elif provider == 'cricket':
            cell.fill.fore_color.rgb = RGBColor(205, 235, 222)
        if chart[provider][key].price_change == 'yes':
            run.font.color.rgb = RGBColor(255, 0, 0)
        else:
            run.font.color.rgb = RGBColor(109, 110, 113)
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(191, 191, 191)
        run.text = ' '

def remove_duplicate_device(chart):
    for a, b in itertools.combinations(chart['master'], 2):
        if a[0] == b[0]:
            if int(a[1]) < int(b[1]):
                if b in chart['master']:
                    chart['master'].remove(b)
            elif int(a[1]) > int(b[1]):
                if a in chart['master']:
                    chart['master'].remove(a)

def make_dict_key(chart, price_dict, provider):
    chart[provider] = {}
    for price in price_dict[provider]:
        for z in range(len(chart['master'])):
            if [fix_spacing(price.device), price.storage] == chart['master'][z]:
                chart[provider].update({z: price})
    return chart[provider]

def limit_chart_size(chart):
    if len(chart['master']) > 25:
        print('chart trimmed:', len(chart['master']) - 25, 'entry/ies removed')
        chart['master'] = chart['master'][:25]
    return chart['master']

def generate_Device_Slides(prs):
    # device chart 1
    slide1 = prs.slides[1]
    shapes = slide1.shapes
    rows = len(comp_chart_1['master']) + 1
    cols = 15
    top = Inches(1.75)
    left = Inches(0.65)
    width = Inches(10.75)
    height = Inches(0)
    table = shapes.add_table(rows, cols, left, top, width, height).table
    table.rows[0].height = Inches(0.5)
    table.columns[0].width = Inches(2.2)
    # color and titles by column
    for y in range(cols):
        cell = table.cell(0, y)
        fill = cell.fill
        tf = cell.text_frame
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        run = p.add_run()
        if y == 0 or y == 1 or y == 2 or y == 3:
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 0, 0)
            run.font.name = 'Ariel'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            if y == 0:
                run.text = "Device"
            if y == 1:
                run.text = "Monthly (24-Mo.)"
            if y == 2:
                run.text = "Retail Price"
            if y == 3:
                run.text = "2-yr Price"
        elif y == 4 or y == 5 or y == 6:
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 112, 192)
            run.font.name = 'Ariel'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            if y == 4:
                run.text = 'Monthly (24-Mo.)'
            if y == 5:
                run.text = 'Retail Price'
            if y == 6:
                run.text = 'Money Down'
        elif y == 7 or y == 8 or y == 9:
            fill.solid()
            fill.fore_color.rgb = RGBColor(210, 102, 159)
            run.font.name = 'Ariel'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            if y == 7:
                run.text = 'Monthly (24-Mo.)'
            if y == 8:
                run.text = 'Retail Price'
            if y == 9:
                run.text = 'Money Down'
        elif y == 10 or y == 11 or y == 12:
            fill.solid()
            fill.fore_color.rgb = RGBColor(74, 154, 77)
            run.font.name = 'Ariel'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            if y == 10:
                run.text = 'Monthly (18-Mo.)'
            if y == 11:
                run.text = 'Retail Price'
            if y == 12:
                run.text = 'Money Down'
        else:
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)
            run.font.name = 'Ariel'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            if y == 13:
                run.text = 'Monthly (24-Mo.)'
            if y == 14:
                run.text = 'Retail Price'
        # by row
        for x in range(len(comp_chart_1['master'])):
            cell = table.cell(x+1, y)
            cell.margin_bottom = 0
            cell.margin_top = 0
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            run = p.add_run()
            if y == 0:
                run.text = format_device_names(comp_chart_1['master'][x])
                run.font.name = 'Ariel'
                run.font.bold = True
                run.font.color.rgb = RGBColor(109, 110, 113)
                p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                run.font.size = Pt(11)
            elif y == 1:
                populate_cell(x, 'verizon', comp_chart_1, cell, run, y)
                run.font.size = Pt(11)
            elif y == 2:
                populate_cell(x, 'verizon', comp_chart_1, cell, run, y)
            elif y == 3:
                populate_cell_contract_ufc(x, 'verizon', comp_chart_1, cell, run)
            elif y == 4:
                populate_cell(x, 'att', comp_chart_1, cell, run, y)
            elif y == 5:
                populate_cell(x, 'att', comp_chart_1, cell, run, y)
            elif y == 6:
                populate_cell(x, 'att', comp_chart_1, cell, run, y)
            elif y == 7:
                populate_cell(x, 'tmobile', comp_chart_1, cell, run, y)
            elif y == 8:
                populate_cell(x, 'tmobile', comp_chart_1, cell, run, y)
            elif y == 9:
                populate_cell(x, 'tmobile', comp_chart_1, cell, run, y)
            elif y == 10:
                populate_cell(x, 'sprint', comp_chart_1, cell, run, y)
            elif y == 11:
                populate_cell(x, 'sprint', comp_chart_1, cell, run, y)
            elif y == 12:
                populate_cell(x, 'sprint', comp_chart_1, cell, run, y)
            elif y == 13:
                populate_cell(x, 'xfinity', comp_chart_1, cell, run, y)
            elif y == 14:
                populate_cell(x, 'xfinity', comp_chart_1, cell, run, y)
            else:
                run.text = ' '
                run.font.size = Pt(9.5)
            run.font.size = Pt(9.5)
            run.font.name = 'Ariel'

    # device chart 2
    slide2 = prs.slides[2]
    shapes = slide2.shapes
    rows = len(tab_chart_2['master']) + 1
    cols = 13
    top = Inches(1.75)
    left = Inches(0.65)
    width = Inches(10.6)
    height = Inches(1)
    table = shapes.add_table(rows, cols, left, top, width, height).table
    table.rows[0].height = Inches(0.5)
    table.columns[0].width = Inches(2.2)
    for y in range(cols):
        cell = table.cell(0, y)
        fill = cell.fill
        tf = cell.text_frame
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        run = p.add_run()
        if y == 0 or y == 1 or y == 2 or y == 3:
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 0, 0)
            run.font.name = 'Ariel'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            if y == 0:
                run.text = "Device"
            if y == 1:
                run.text = "Monthly (24-Mo.)"
            if y == 2:
                run.text = "Retail Price"
            if y == 3:
                run.text = "2-yr Price"
        elif y == 4 or y == 5 or y == 6:
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 112, 192)
            run.font.name = 'Ariel'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            if y == 4:
                run.text = 'Monthly (24-Mo.)'
            if y == 5:
                run.text = 'Retail Price'
            if y == 6:
                run.text = '24-Mo. Contract (UFC)'
        elif y == 7 or y == 8 or y == 9:
            fill.solid()
            fill.fore_color.rgb = RGBColor(210, 102, 159)
            run.font.name = 'Ariel'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            if y == 7:
                run.text = 'Monthly (24-Mo.)'
            if y == 8:
                run.text = 'Retail Price'
            if y == 9:
                run.text = 'Money Down'
        else:
            fill.solid()
            fill.fore_color.rgb = RGBColor(74, 154, 77)
            run.font.name = 'Ariel'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            if y == 10:
                run.text = 'Monthly (18-Mo.)'
            if y == 11:
                run.text = 'Retail Price'
            if y == 12:
                run.text = 'Money Down'
        # by row
        for x in range(len(tab_chart_2['master'])):
            cell = table.cell(x+1, y)
            cell.margin_bottom = 0
            cell.margin_top = 0
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            run = p.add_run()
            if y == 0:
                run.text = format_device_names(tab_chart_2['master'][x])
                run.font.name = 'Ariel'
                run.font.bold = True
                run.font.color.rgb = RGBColor(109, 110, 113)
                p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            elif y == 1:
                populate_cell(x, 'verizon', tab_chart_2, cell, run, y)
            elif y == 2:
                populate_cell(x, 'verizon', tab_chart_2, cell, run, y)
            elif y == 3:
                populate_cell_contract_ufc(x, 'verizon', tab_chart_2, cell, run)
            elif y == 4:
                populate_cell(x, 'att', tab_chart_2, cell, run, y)
            elif y == 5:
                populate_cell(x, 'att', tab_chart_2, cell, run, y)
            elif y == 6:
                populate_cell_contract_ufc(x, 'att', tab_chart_2, cell, run)
            elif y == 7:
                populate_cell(x, 'tmobile', tab_chart_2, cell, run, y)
            elif y == 8:
                populate_cell(x, 'tmobile', tab_chart_2, cell, run, y)
            elif y == 9:
                populate_cell(x, 'tmobile', tab_chart_2, cell, run, y)
            elif y == 10:
                populate_cell(x, 'sprint', tab_chart_2, cell, run, y)
            elif y == 11:
                populate_cell(x, 'sprint', tab_chart_2, cell, run, y)
            elif y == 12:
                populate_cell(x, 'sprint', tab_chart_2, cell, run, y)
            else:
                run.text = ' '
            run.font.size = Pt(11)
            run.font.name = 'Ariel'

    # device chart 3
    slide3 = prs.slides[3]
    shapes = slide3.shapes
    rows = len(sub10_chart_3['master']) + 1
    cols = 15
    top = Inches(1.75)
    left = Inches(0.65)
    width = Inches(10.5)
    height = Inches(1)
    table = shapes.add_table(rows, cols, left, top, width, height).table
    table.rows[0].height = Inches(0.2)
    table.columns[0].width = Inches(2.2)
    for y in range(cols):
        cell = table.cell(0, y)
        fill = cell.fill
        tf = cell.text_frame
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        run = p.add_run()
        if y == 0 or y == 1 or y == 2 or y == 3:
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 0, 0)
            run.font.name = 'Ariel'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            if y == 0:
                run.text = "Device"
            if y == 1:
                run.text = "Monthly (24-Mo.)"
            if y == 2:
                run.text = "Retail Price"
            if y == 3:
                run.text = "2-yr Price"
        elif y == 4 or y == 5 or y == 6:
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 112, 192)
            run.font.name = 'Ariel'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            if y == 4:
                run.text = 'Monthly (24-Mo.)'
            if y == 5:
                run.text = 'Retail Price'
            if y == 6:
                run.text = 'Money Down'
        elif y == 7 or y == 8 or y == 9:
            fill.solid()
            fill.fore_color.rgb = RGBColor(210, 102, 159)
            run.font.name = 'Ariel'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            if y == 7:
                run.text = 'Monthly (24-Mo.)'
            if y == 8:
                run.text = 'Retail Price'
            if y == 9:
                run.text = 'Money Down'
        elif y == 10 or y == 11 or y == 12:
            fill.solid()
            fill.fore_color.rgb = RGBColor(74, 154, 77)
            run.font.name = 'Ariel'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            if y == 10:
                run.text = 'Monthly (18-Mo.)'
            if y == 11:
                run.text = 'Retail Price'
            if y == 12:
                run.text = 'Money Down'
        else:
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 0, 0)
            run.font.name = 'Ariel'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            if y == 13:
                run.text = 'Monthly (24-Mo.)'
            if y == 14:
                run.text = 'Retail Price'
        for x in range(len(sub10_chart_3['master'])):
            cell = table.cell(x+1, y)
            cell.margin_bottom = 0
            cell.margin_top = 0
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            run = p.add_run()
            if y == 0:
                run.text = format_device_names(sub10_chart_3['master'][x])
                run.font.name = 'Ariel'
                run.font.bold = True
                run.font.color.rgb = RGBColor(109, 110, 113)
                p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            elif y == 1:
                populate_cell(x, 'verizon', sub10_chart_3, cell, run, y)
            elif y == 2:
                populate_cell(x, 'verizon', sub10_chart_3, cell, run, y)
            elif y == 3:
                populate_cell_contract_ufc(x, 'verizon', sub10_chart_3, cell, run)
            elif y == 4:
                populate_cell(x, 'att', sub10_chart_3, cell, run, y)
            elif y == 5:
                populate_cell(x, 'att', sub10_chart_3, cell, run, y)
            elif y == 6:
                populate_cell(x, 'att', sub10_chart_3, cell, run, y)
            elif y == 7:
                populate_cell(x, 'tmobile', sub10_chart_3, cell, run, y)
            elif y == 8:
                populate_cell(x, 'tmobile', sub10_chart_3, cell, run, y)
            elif y == 9:
                populate_cell(x, 'tmobile', sub10_chart_3, cell, run, y)
            elif y == 10:
                populate_cell(x, 'sprint', sub10_chart_3, cell, run, y)
            elif y == 11:
                populate_cell(x, 'sprint', sub10_chart_3, cell, run, y)
            elif y == 12:
                populate_cell(x, 'sprint', sub10_chart_3, cell, run, y)
            elif y == 13:
                populate_cell(x, 'xfinity', sub10_chart_3, cell, run, y)
            elif y == 14:
                populate_cell(x, 'xfinity', sub10_chart_3, cell, run, y)
            else:
                run.text = ' '
            run.font.size = Pt(10.5)
            run.font.name = 'Ariel'

    # device chart 4
    slide4 = prs.slides[4]
    shapes = slide4.shapes
    rows = len(prepaid_comp_chart_4['master']) + 1
    cols = 5
    top = Inches(1.75)
    left = Inches(0.65)
    width = Inches(11.85)
    height = Inches(1)
    table = shapes.add_table(rows, cols, left, top, width, height).table
    table.rows[0].height = Inches(0.1)
    table.columns[0].width = Inches(2.5)
    for y in range(cols):
        cell = table.cell(0, y)
        fill = cell.fill
        tf = cell.text_frame
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        run = p.add_run()
        if y == 0 or y == 1:
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 0, 0)
            run.font.name = 'Ariel'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            if y == 0:
                run.text = 'Devices'
            if y == 1:
                run.text = 'Retail Price'
        elif y == 2:
            fill.solid()
            fill.fore_color.rgb = RGBColor(0, 112, 192)
            run.font.name = 'Ariel'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.text = 'Retail Price'
        elif y == 3:
            fill.solid()
            fill.fore_color.rgb = RGBColor(244, 110, 55)
            run.font.name = 'Ariel'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.text = 'Retail Price'
        else:
            fill.solid()
            fill.fore_color.rgb = RGBColor(146, 208, 80)
            run.font.name = 'Ariel'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.text = 'Retail Price'
        for x in range(len(prepaid_comp_chart_4['master'])):
            cell = table.cell(x+1, y)
            cell.margin_top = 0
            cell.margin_bottom = 0
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            run = p.add_run()
            if y == 0:
                run.text = format_device_names(prepaid_comp_chart_4['master'][x])
                run.font.name = 'Ariel'
                run.font.bold = True
                run.font.size = Pt(11)
                run.font.color.rgb = RGBColor(109, 110, 113)
                p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            elif y == 1:
                prepaid_populate_cell(x, 'verizon', prepaid_comp_chart_4, cell, run)
            elif y == 2:
                prepaid_populate_cell(x, 'att', prepaid_comp_chart_4, cell, run)
            elif y == 3:
                prepaid_populate_cell(x, 'metropcs', prepaid_comp_chart_4, cell, run)
            elif y == 4:
                prepaid_populate_cell(x, 'cricket', prepaid_comp_chart_4, cell, run)
            else:
                run.text = ' '
        run.font.name = 'Ariel'

    # device chart 5
    slide5 = prs.slides[5]
    shapes = slide5.shapes
    rows = len(cri_metro_chart_5['master']) + 1
    cols = 3
    top = Inches(1.75)
    left = Inches(0.65)
    width = Inches(12)
    height = Inches(1)
    table = shapes.add_table(rows, cols, left, top, width, height).table
    table.rows[0].height = Inches(0.1)
    table.columns[0].width = Inches(4)
    # by column, first row
    for y in range(cols):
        cell = table.cell(0, y)
        fill = cell.fill
        tf = cell.text_frame
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        run = p.add_run()
        if y == 0:
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 0, 0)
            run.font.name = 'Ariel'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.text = 'Device'
        elif y == 1:
            fill.solid()
            fill.fore_color.rgb = RGBColor(244, 110, 55)
            run.font.name = 'Ariel'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.text = 'Retail Price'
        else:
            fill.solid()
            fill.fore_color.rgb = RGBColor(146, 208, 80)
            run.font.name = 'Ariel'
            run.font.size = Pt(10)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.text = 'Retail Price'
        # by row, rows after first
        for x in range(len(cri_metro_chart_5['master'])):
            cell = table.cell(x+1, y)
            cell.margin_top = 0
            cell.margin_bottom = 0
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            run = p.add_run()
            if y == 0:
                run.text = format_device_names(cri_metro_chart_5['master'][x])
                run.font.name = 'Ariel'
                run.font.bold = True
                run.font.size = Pt(11)
                run.font.color.rgb = RGBColor(109, 110, 113)
                p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            elif y == 1:
                prepaid_populate_cell(x, 'metropcs', cri_metro_chart_5, cell, run)
            elif y == 2:
                prepaid_populate_cell(x, 'cricket', cri_metro_chart_5, cell, run)
            else:
                run.text = ' '
        run.font.name = 'Ariel'

# postpaid devices
# make master list (and first column labels) by chart
comp_chart_1 = {}
tab_chart_2 = {}
sub10_chart_3 = {}
postpaid_prices = {}
comp_chart_1['master'] = []
tab_chart_2['master'] = []
sub10_chart_3['master'] = []
for provider in postpaid_providers:
    postpaid_prices[provider] = get_postpaid_device_prices(provider, today)
    for price in postpaid_prices[provider]:
        if price.is_tablet == 'yes':
            if [fix_spacing(price.device), price.storage] not in tab_chart_2['master']:
                tab_chart_2['master'].append([fix_spacing(price.device), price.storage])
        elif price.monthly_price > 15 and price.is_flipphone == 'no':
            if [fix_spacing(price.device), price.storage] not in comp_chart_1['master']:
                comp_chart_1['master'].append([fix_spacing(price.device), price.storage])
        elif price.is_flipphone == 'no':
            if [fix_spacing(price.device), price.storage] not in sub10_chart_3['master'] and [fix_spacing(price.device), price.storage] not in comp_chart_1['master']:
                sub10_chart_3['master'].append([fix_spacing(price.device), price.storage])

# if the same device is in the list more than once, remove the larger storage size
remove_duplicate_device(comp_chart_1)
remove_duplicate_device(tab_chart_2)
remove_duplicate_device(sub10_chart_3)

# make dictionary of prices with row number as key
for provider in postpaid_providers:
    make_dict_key(comp_chart_1, postpaid_prices, provider)
    make_dict_key(tab_chart_2, postpaid_prices, provider)
    make_dict_key(sub10_chart_3, postpaid_prices, provider)

# prepaid devices
# make master list (and first column labels) by chart
prepaid_comp_chart_4 = {}
cri_metro_chart_5 = {}
prepaid_comp_chart_4['master'] = []
cri_metro_chart_5['master'] = []
prepaid_prices = {}
for provider in prepaid_providers:
    prepaid_prices[provider] = get_prepaid_device_prices(provider, today)
    chart_5_limit = 0
    for price in prepaid_prices[provider]:
        if (price.provider == 'att' or price.provider == 'verizon') and price.price >= 300:
            if [fix_spacing(price.device), price.storage] not in prepaid_comp_chart_4['master']:
                prepaid_comp_chart_4['master'].append([fix_spacing(price.device), price.storage])
        elif price.price >= 300:
            if [fix_spacing(price.device), price.storage] not in prepaid_comp_chart_4['master']:
                prepaid_comp_chart_4['master'].append([fix_spacing(price.device), price.storage])
        elif price.device.find('iPhone') != -1 or price.device.find('Galaxy S') != -1:
            if [fix_spacing(price.device), price.storage] not in prepaid_comp_chart_4['master']:
                prepaid_comp_chart_4['master'].append([fix_spacing(price.device), price.storage])
        else:
            if (price.provider == 'metropcs' or price.provider == 'cricket') \
                    and [fix_spacing(price.device), price.storage] not in cri_metro_chart_5['master']\
                    and price.is_flipphone == 'no':
                if chart_5_limit > 11:
                    continue
                elif chart_5_limit <= 11:
                    cri_metro_chart_5['master'].append([fix_spacing(price.device), price.storage])
                    chart_5_limit += 1

# if the same device is in the list more than once, remove the larger storage size
remove_duplicate_device(prepaid_comp_chart_4)
remove_duplicate_device(cri_metro_chart_5)

# make dictionary of prices with row number as key
for provider in prepaid_providers:
    make_dict_key(prepaid_comp_chart_4, prepaid_prices, provider)
    make_dict_key(cri_metro_chart_5, prepaid_prices, provider)

# max out the device charts at 25
comp_chart_1['master'] = limit_chart_size(comp_chart_1)
tab_chart_2['master'] = limit_chart_size(tab_chart_2)
sub10_chart_3['master'] = limit_chart_size(sub10_chart_3)
prepaid_comp_chart_4['master'] = limit_chart_size(prepaid_comp_chart_4)
# cri_metro_chart_5 is already limited to 22 (11 each provider)

# # print number of rows, device names for each chart
# print('comp chart: ', len(comp_chart_1['master']), comp_chart_1['master'])
# print('tablets chart: ', len(tab_chart_2['master']), tab_chart_2['master'])
# print('sub10 chart: ', len(sub10_chart_3['master']), sub10_chart_3['master'])
# print('pre comp chart: ', len(prepaid_comp_chart_4['master']), prepaid_comp_chart_4['master'])
# print('cri/met chart: ', len(cri_metro_chart_5['master']), cri_metro_chart_5['master'])
